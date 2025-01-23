from pptx import Presentation

def create_ppt(content, output_path):
    try:
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        title = slide.shapes.title
        content_box = slide.placeholders[1]

        title.text = "Scraped Content"
        content_box.text = "\n".join(content)

        prs.save(output_path)
        print(f"Presentation saved to {output_path}")
    except Exception as e:
        print(f"Error creating PowerPoint: {e}")

if __name__ == "__main__":
    sample_content = ["Example Slide 1", "More data here"]
    create_ppt(sample_content, "output.pptx")
